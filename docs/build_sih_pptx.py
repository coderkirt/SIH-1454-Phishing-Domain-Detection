"""Build the official 6-slide SIH 2026 idea PPT from the SIH template + live codebase."""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs" / "_sih_template.pptx"
OUT = ROOT / "PhishShield_SIH_2026_Final_Presentation.pptx"

NAVY = RGBColor(0x1F, 0x49, 0x7D)
ORANGE = RGBColor(0xF7, 0x96, 0x46)
GREEN = RGBColor(0x3D, 0x8B, 0x5F)
RED = RGBColor(0xC0, 0x50, 0x4D)
BLUE = RGBColor(0x4F, 0x81, 0xBD)
TEAL = RGBColor(0x4B, 0xAC, 0xC6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
AMBER = RGBColor(0xD9, 0x77, 0x06)
CARD = RGBColor(0xEE, 0xF4, 0xFA)


def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def set_runs(shape, lines, size=14, bold=False, color=DARK):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.clear()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(4)


def box(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line or fill
    sh.line.width = Pt(0.75)
    return sh


def label(shape, text, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    try:
        tf.paragraphs[0].space_before = Pt(0)
    except Exception:
        pass


def add_text(slide, l, t, w, h, text, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return sh


def hide_body(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if "Proposed Solution" in t or "Technologies to be used" in t or "Analysis of the feasibility" in t:
            shape.left = Inches(13.2)
            shape.width = Inches(0.1)
            shape.height = Inches(0.1)
        if "Details / Links" in t or "Potential impact" in t:
            shape.left = Inches(13.2)
            shape.width = Inches(0.1)
            shape.height = Inches(0.1)
        if shape.text_frame.text.strip() == "Your Team Name":
            set_runs(shape, ["[TEAM NAME]"], size=12, bold=True, color=NAVY)


def fill_title(prs):
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "Problem Statement ID" in text:
            set_runs(
                shape,
                [
                    "Problem Statement ID – SIH1454",
                    "Problem Statement Title – Create an intelligent system using AI/ML to detect phishing domains which imitate look and feel of genuine domains",
                    "Theme – Blockchain & Cybersecurity",
                    "PS Category – Software",
                    "Team ID – [TEAM ID]",
                    "Team Name (Registered on portal) – [TEAM NAME]",
                ],
                size=16,
                bold=True,
                color=DARK,
            )
        elif "TITLE PAGE" in text:
            set_runs(shape, ["TITLE PAGE"], size=28, bold=True, color=NAVY)


def slide2(prs):
    slide = prs.slides[1]
    hide_body(slide)
    add_text(slide, 0.4, 1.05, 12.4, 0.28, "IDEA TITLE  |  PhishShield AI   —  Analyze before you trust", 16, True, NAVY)
    add_text(slide, 0.4, 1.32, 12.4, 0.28, "Proposed Solution  ·  How it addresses the problem  ·  Innovation and uniqueness", 11, False, MUTED)

    # input row
    for i, (txt, col) in enumerate([
        ("URL / domain", BLUE),
        ("Pasted email / WA / SMS", TEAL),
        ("Webpage (extension)", NAVY),
        ("QR image", GREEN),
    ]):
        b = box(slide, 0.4 + i * 3.2, 1.65, 3.05, 0.42, col)
        label(b, txt, 11, True, WHITE)

    arrow = box(slide, 5.15, 2.12, 3.0, 0.32, ORANGE)
    label(arrow, "↓  POST /api/v1/analyze/content", 10, True, WHITE)

    engines = [
        ("URL / domain\nanalysis", BLUE),
        ("Psychological\nthreat (EN+Hinglish)", RED),
        ("Brand +\nsender signals", ORANGE),
        ("Threat intel\nGSB / OpenPhish", NAVY),
        ("Community\nreports (signal)", TEAL),
    ]
    for i, (txt, col) in enumerate(engines):
        b = box(slide, 0.35 + i * 2.58, 2.52, 2.48, 0.72, col)
        label(b, txt, 10, True, WHITE)

    fuse = box(slide, 3.6, 3.32, 6.1, 0.38, DARK)
    label(fuse, "RISK FUSION ENGINE  →  score 0–100  ·  LOW / MED / HIGH / CRITICAL", 12, True, WHITE)

    s = box(slide, 1.4, 3.78, 4.8, 0.42, GREEN)
    label(s, "SIMPLE VIEW  —  plain-English warning", 12, True, WHITE)
    t = box(slide, 7.0, 3.78, 4.8, 0.42, NAVY)
    label(t, "TECHNICAL VIEW  —  evidence + tags", 12, True, WHITE)

    add_text(slide, 0.4, 4.28, 12.5, 0.28, "How it addresses SIH1454 (NTRO): imitate-genuine-domain detection + usable warning, not jargon.", 12, True, DARK)

    cards = [
        ("1 DETECT", "Lookalike domains, short links, IP URLs, new age, intel flags"),
        ("2 UNDERSTAND", "Urgency / fear / coercion / Hinglish in pasted messages"),
        ("3 EXPLAIN", "Simple View: do not enter OTP. Technical View: why."),
        ("4 PROTECT", "Browser overlay on HIGH/CRITICAL. Same API as dashboard."),
        ("5 LEARN", "Authenticated Mark as Scam/Risky/Safe — signal, not proof."),
    ]
    for i, (h, btxt) in enumerate(cards):
        c = box(slide, 0.35 + i * 2.58, 4.58, 2.48, 1.15, CARD, BLUE)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = NAVY
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = btxt
        r2.font.size = Pt(9)
        r2.font.color.rgb = MUTED
        r2.font.name = "Calibri"

    add_text(
        slide,
        0.4,
        5.78,
        12.5,
        0.55,
        "Honesty: pasted content (not Gmail/WhatsApp inbox). Screenshot OCR = PARTIAL (Tesseract). Native phone SMS = PLANNED. "
        "‘ML’ fusion weight currently reuses URL risk — trained CNN/classifier = PLANNED. Confidence = signal agreement, not accuracy.",
        10,
        False,
        AMBER,
    )


def slide3(prs):
    slide = prs.slides[2]
    hide_body(slide)
    add_text(slide, 0.35, 1.02, 12.6, 0.26, "TECHNICAL APPROACH  ·  Technologies  ·  Methodology  ·  Working prototype", 14, True, NAVY)

    stacks = [
        ("FRONTEND", "React 19  ·  Vite 8\nTailwind  ·  Axios  ·  Recharts", BLUE),
        ("BACKEND", "Python  ·  FastAPI\nPydantic  ·  Uvicorn", NAVY),
        ("DATABASE", "SQLite (live)\nPostgreSQL = PLANNED", TEAL),
        ("DETECTION", "URLChecker  ·  Psychology 2.0\nRisk fusion  ·  OpenCV QR", ORANGE),
        ("EXTENSION", "Manifest V3\nservice worker + overlay", GREEN),
        ("SECURITY", "bcrypt  ·  JWT  ·  SSRF guards\nparameterized SQL", RED),
    ]
    for i, (h, body, col) in enumerate(stacks):
        x = 0.35 + (i % 6) * 2.15
        b = box(slide, x, 1.32, 2.08, 1.15, col)
        tf = b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(10)
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE
        r2.font.name = "Calibri"

    add_text(slide, 0.35, 2.55, 12.6, 0.24, "Architecture (actual runtime — not Flask, not a trained CNN)", 12, True, DARK)

    flow = [
        ("INPUT", "URL · paste · QR\nextension tab"),
        ("NORMALIZE", "extract URLs\nunwrap short links"),
        ("ANALYZE", "domain · intel\npsych · brand · sender"),
        ("FUSE", "weighted signals\npsych cannot solo-CRITICAL"),
        ("EXPLAIN", "Simple + Technical\nviews"),
        ("ACT", "warn · report\nstore SQLite"),
    ]
    for i, (h, btxt) in enumerate(flow):
        c = box(slide, 0.35 + i * 2.15, 2.82, 2.05, 0.95, CARD, NAVY)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = btxt
        r2.font.size = Pt(9)
        r2.font.color.rgb = MUTED

    add_text(slide, 0.35, 3.85, 12.6, 0.22, "Process: Extract → Analyze → Score → Explain → Warn / Report     |     Prototype: React :5173  +  FastAPI :8000  +  unpacked Chrome extension", 11, True, DARK)

    add_text(
        slide,
        0.35,
        4.12,
        12.6,
        1.95,
        "NOT in this repo: Flask  ·  live PostgreSQL  ·  TensorFlow/Keras CNN  ·  sklearn model file  ·  Gmail/Outlook OAuth  ·  native Android SMS.\n"
        "Intel: Google Safe Browsing if API key set, else OpenPhish. QR: OpenCV. Screenshot: QR always; OCR only if Tesseract is installed.\n"
        "Key files: backend/app/services/content_analyzer.py  ·  risk_engine.py  ·  url_checker.py  ·  psychology.py  ·  extension/  ·  frontend/src/pages/Scanner.jsx",
        12,
        False,
        MUTED,
    )


def slide4(prs):
    slide = prs.slides[3]
    hide_body(slide)
    add_text(slide, 0.35, 1.02, 12.6, 0.28, "FEASIBILITY AND VIABILITY  ·  Analysis  ·  Challenges  ·  Strategies", 14, True, NAVY)

    cols = [
        ("TECHNICAL", GREEN, "Working FastAPI+React prototype\npytest coverage on APIs\nModular services (swap SQLite→Postgres later)\nOpen-source stack, no paid licence required\nQR/OpenCV already in requirements"),
        ("OPERATIONAL", BLUE, "Paste URL/message in Scanner tabs\nExtension scans current tab + visible links\nPlain-language Simple View for non-tech users\nDashboard history / stats from SQLite\nNo need to read private inboxes"),
        ("SECURITY", NAVY, "IMPLEMENTED: bcrypt, JWT, ?-SQL,\nSSRF block (private IP/metadata),\nreport 24h/20-cap, 5 MB upload cap\nPARTIAL: CORS allow-all (dev)\nPLANNED: global API rate-limit, cert-chain"),
    ]
    for i, (h, col, body) in enumerate(cols):
        b = box(slide, 0.35 + i * 4.3, 1.38, 4.15, 2.15, CARD, col)
        tf = b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = col
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK

    add_text(slide, 0.35, 3.62, 12.6, 0.24, "Challenge  →  mitigation (implemented unless marked planned)", 12, True, DARK)
    rows = [
        ("False positives from marketing language", "Psych cannot raise HIGH/CRITICAL alone (risk_engine.py)"),
        ("New phishing domains / shorteners", "RDAP age + unwrap redirects + GSB/OpenPhish + brand distance"),
        ("Community report abuse", "Login required + 1 URL/24h + 20 reports/day"),
        ("Privacy / inbox access", "User-pasted content only; no OTP/password scrape"),
        ("SQLite at multi-server scale", "PLANNED: PostgreSQL; MVP is single-node demo"),
        ("No labelled CNN yet", "PLANNED trained model; today explainable fusion + reserved ml_confidence"),
    ]
    y = 3.88
    for ch, mit in rows:
        left = box(slide, 0.35, y, 5.9, 0.32, RGBColor(0xFD, 0xE8, 0xE8), RED)
        label(left, ch, 10, True, DARK, PP_ALIGN.LEFT)
        right = box(slide, 6.4, y, 6.5, 0.32, RGBColor(0xE8, 0xF7, 0xEF), GREEN)
        label(right, mit, 10, False, DARK, PP_ALIGN.LEFT)
        y += 0.36


def slide5(prs):
    slide = prs.slides[4]
    hide_body(slide)
    add_text(slide, 0.35, 1.02, 12.6, 0.26, "IMPACT AND BENEFITS  ·  Target audience  ·  Social / economic value", 14, True, NAVY)

    users = ["Elderly users", "First-time UPI / banking", "Students", "Everyday smartphone users", "Non-technical staff", "SOC-curious learners"]
    for i, u in enumerate(users):
        b = box(slide, 0.35 + (i % 6) * 2.15, 1.35, 2.08, 0.38, BLUE if i % 2 == 0 else NAVY)
        label(b, u, 10, True, WHITE)

    add_text(slide, 0.35, 1.82, 6.2, 0.24, "BEFORE", 12, True, RED)
    add_text(slide, 6.9, 1.82, 6.2, 0.24, "WITH PHISHSHIELD", 12, True, GREEN)

    before = box(slide, 0.35, 2.08, 6.15, 2.15, RGBColor(0xFD, 0xED, 0xED), RED)
    label(
        before,
        "Suspicious SMS / WhatsApp link\n→ User sees padlock / jargon\n→ “SSL mismatch” is meaningless\n→ Enters OTP on fake bank page\n→ Credential theft",
        13,
        False,
        DARK,
        PP_ALIGN.LEFT,
    )
    after = box(slide, 6.85, 2.08, 6.15, 2.15, RGBColor(0xEC, 0xF8, 0xF1), GREEN)
    label(
        after,
        "Same link pasted or opened in browser\n→ Fusion score + threat tags\n→ “This may be pretending to be SBI/PayPal.\n    Do not enter password or OTP.”\n→ Go back  or  Continue anyway (conscious)",
        13,
        False,
        DARK,
        PP_ALIGN.LEFT,
    )

    add_text(slide, 0.35, 4.32, 12.6, 0.24, "Benefits: fewer credential-harvest victims  ·  explainable warnings  ·  community signal  ·  one engine for web + extension", 12, True, NAVY)
    add_text(slide, 0.35, 4.58, 12.6, 0.22, "Proposed evaluation metrics  (NOT claimed results — no published F1 in this repository)", 11, True, AMBER)

    kpis = [
        "Detection recall on labelled lookalike domains",
        "False-positive rate on trusted banks",
        "Median analysis latency",
        "Share of HIGH/CRITICAL with usable Simple View",
        "Community report precision after review",
        "Channel coverage: URL / paste / QR / extension",
    ]
    for i, k in enumerate(kpis):
        b = box(slide, 0.35 + (i % 3) * 4.3, 4.88 + (i // 3) * 0.55, 4.15, 0.48, CARD, ORANGE)
        label(b, k, 11, False, DARK, PP_ALIGN.LEFT)


def slide6(prs):
    slide = prs.slides[5]
    hide_body(slide)
    add_text(slide, 0.35, 1.02, 12.6, 0.28, "RESEARCH AND REFERENCES  ·  Sources actually used by this implementation", 14, True, NAVY)

    refs = [
        ("SIH1454 / NTRO", "Create an intelligent system using AI/ML to detect phishing domains which imitate look and feel of genuine domains. Theme: Blockchain & Cybersecurity. Category: Software."),
        ("Google Safe Browsing API v4", "https://developers.google.com/safe-browsing — optional live IOC lookup (threatMatches:find)."),
        ("OpenPhish", "https://openphish.com — public phishing URL feed used when no Google key is configured."),
        ("RDAP (IETF RFC 7480+)", "https://rdap.org — domain registration age for newly registered lookalikes."),
        ("OWASP", "Phishing / unvalidated redirects / SSRF cheat sheets — informed SSRF blocks in ad_signals.py / ssrf.py."),
        ("FastAPI", "https://fastapi.tiangolo.com — actual API framework (not Flask)."),
        ("React + Vite", "https://react.dev  ·  https://vite.dev — dashboard SPA."),
        ("Chrome Extensions MV3", "https://developer.chrome.com/docs/extensions — service worker, content scripts, host permissions."),
        ("SQLite", "https://www.sqlite.org — live database; PostgreSQL remains a scale-up plan."),
        ("OpenCV QR", "opencv-python-headless in requirements.txt — QR decode from uploaded images."),
    ]
    y = 1.35
    for title, body in refs:
        add_text(slide, 0.45, y, 3.3, 0.42, title, 11, True, NAVY)
        add_text(slide, 3.8, y, 9.1, 0.42, body, 10, False, DARK)
        y += 0.46
    add_text(
        slide,
        0.45,
        6.05,
        12.4,
        0.35,
        "No fabricated papers or accuracy numbers. Fill Team ID / Team Name on slide 1 from the SIH portal before upload.",
        11,
        True,
        AMBER,
    )


def main():
    prs = Presentation(str(TEMPLATE))
    # drop instruction slide 7
    delete_slide(prs, 6)
    fill_title(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    assert len(prs.slides) == 6, len(prs.slides)
    prs.save(str(OUT))
    print("Wrote", OUT, "slides", len(prs.slides))


if __name__ == "__main__":
    main()
